"""All RAG bootcamp generators in one file (calls discrete builders)."""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from helpers import (
    new_doc, add_title, add_subtitle, add_h1, add_h2, add_h3, add_body,
    add_bullet, add_numbered, add_code, add_table, add_divider, add_page_break,
)
from bootcamp_content import (
    WEEKS, slug, START_DATE, DURATION_WEEKS, TUITION, PREREQ,
)
from class_notes_deep import DEEP
from quiz_bank import QUIZZES
from projects_bank import PROJECTS

ROOT = "/Users/pkr465/work/AI-Course-Material/LLM_and_Enterprise_RAG_Bootcamp"
COURSE_NAME = "LLM and Enterprise RAG Bootcamp"
INSTRUCTOR = "Pavan R"
LETTERS = ["A", "B", "C", "D"]


# ------------------ catalogue + syllabus ------------------

def build_catalogue():
    doc = new_doc()
    add_title(doc, "Proxiant Academy", size=24)
    add_subtitle(doc, COURSE_NAME, size=16)
    add_subtitle(doc, "Course Catalogue", size=12)
    add_divider(doc)

    add_body(doc, f"Start date: {START_DATE}")
    add_body(doc, f"Duration: {DURATION_WEEKS} weeks")
    add_body(doc, "Format: Hybrid (in-person at Fremont campus, live Zoom, or both)")
    add_body(doc, f"Tuition: {TUITION}")
    add_body(doc, "Schedule: Saturdays 11 AM PST through 5 PM PST (main session). Monday and Wednesday 7 to 10 PM PST (labs). Tuesday 8:30 to 10 AM PST (summary and quiz).")
    add_body(doc, "")

    add_h1(doc, "Who This Is For")
    add_body(doc,
        "This bootcamp is built for engineers and ML practitioners who are "
        "comfortable with Python and basic ML and want production-grade depth "
        "in LLMs, embeddings, and retrieval-augmented generation. Teams of "
        "4-6 work in dedicated rooms with full multimedia setup. Compute "
        "access on the Proxiant Datacenter is included and may be extended "
        "by four weeks for project completion.")

    add_h1(doc, "Learning Outcomes")
    outcomes = [
        "Build production RAG systems with sparse + dense retrieval, Matryoshka pruning, ColBERT, and cross-encoder fusion.",
        "Implement four chunking strategies and hierarchical retrieval.",
        "Add multimodal capabilities using CLIP, BLIP, and BLIP-2.",
        "Optimize prompts programmatically with DSPy, GEPA, TextGrad, and ORPO.",
        "Build graph-based retrieval with RAPTOR, GraphRAG, and LightRAG.",
        "Engineer guardrail pipelines and NLI-based grounding.",
        "Apply derivative artifacts, HyDE, semantic caching, and hard-negative mining.",
        "Build agentic RAG and Text-to-SQL with MCP for tool integration.",
        "Fine-tune embedders and base models with PEFT, LoRA, qLoRA, and RLVR.",
        "Evaluate RAG honestly using Ragas, LLM-as-judge, RGB-style robustness, and verification asymmetry.",
    ]
    for o in outcomes:
        add_bullet(doc, o)

    add_h1(doc, "Skills You Will Learn")
    skills = [
        ("Retrieval engineering", "BM25, SPLADE, dense, Matryoshka, ColBERT, cross-encoder, RRF."),
        ("Chunking", "Fixed, semantic, late, contextual, hierarchical."),
        ("Graph-based retrieval", "RAPTOR trees, GraphRAG (triplets, ontologies, communities), LightRAG."),
        ("Multimodal", "Vision transformers, CLIP, BLIP/BLIP-2, Q-Former, multimodal RAG patterns."),
        ("Prompt optimization", "CO-STAR, metaprompting, DSPy (COPRO, MIPRO), GEPA, TextGrad, ORPO."),
        ("Production engineering", "Guardrails, NLI grounding, semantic caching, P95 latency budgets."),
        ("Agentic RAG and SQL", "MCP, Google ADK, agent lifecycles, Text-to-SQL, CTE libraries."),
        ("Fine-tuning", "Chinchilla, LoRA, qLoRA, RLHF, RLVR, GRPO."),
        ("Evaluation", "Retrieval metrics, Ragas, LLM-as-judge bias control, RGB benchmark."),
    ]
    add_table(doc, ["Area", "Coverage"], skills, col_widths=[1.8, 5.0])

    add_h1(doc, "What Is Included")
    included = [
        "12 main lecture sessions (Saturdays, full day) with morning theory and afternoon labs.",
        "Team rooms with full multimedia setup for in-person teams.",
        "24 guided lab sessions (Monday and Wednesday evenings) with working code.",
        "12 weekly quizzes (Tuesday mornings) covering theory and labs.",
        "12 weekly team projects with peer review.",
        "12 weekly research paper reading guides with discussion questions.",
        "Access to the Proxiant Datacenter: 20+ GPU servers, 40+ NVIDIA GPUs (RTX PRO 6000 Blackwell, RTX 5090, RTX 4090), high-speed networking.",
        "Four-week compute extension after course end for project completion.",
        "Full session recordings on the workshop portal, available indefinitely.",
        "Capstone presentation slot and peer review.",
        "PCAP-RAG certificate on passing the final exam.",
    ]
    for x in included:
        add_bullet(doc, x)

    add_h1(doc, "Prerequisites")
    add_body(doc, PREREQ)

    add_h1(doc, "Faculty")
    add_body(doc, f"Lead instructor: {INSTRUCTOR}. Teaching assistants provide 1-on-1 sessions and lab support throughout.")

    add_h1(doc, "Tuition and Registration")
    add_body(doc, TUITION)
    add_body(doc,
        "Register at proxiant.com/rag-bootcamp or email info@proxiant.com. "
        f"Day 1 ({START_DATE}) is open to all interested participants free "
        "of charge as an introduction to the program.")

    add_h1(doc, "Twelve-Week Map")
    rows = [(f"Week {w['num']}", w["date"].split(", ")[1], w["title"], w["tagline"]) for w in WEEKS]
    add_table(doc, ["Week", "Date", "Title", "Theme"], rows, col_widths=[0.7, 1.4, 2.2, 2.6])

    out = os.path.join(ROOT, "Catalogue", "LLM_Enterprise_RAG_Course_Catalogue.docx")
    doc.save(out)
    print(f"WROTE {out}")


def build_syllabus():
    doc = new_doc()
    add_title(doc, "Proxiant Academy", size=24)
    add_subtitle(doc, COURSE_NAME, size=16)
    add_subtitle(doc, "Detailed Twelve-Week Syllabus", size=12)
    add_divider(doc)

    add_h1(doc, "Course Information")
    info = [
        ("Course title", COURSE_NAME),
        ("Start date", START_DATE),
        ("Duration", f"{DURATION_WEEKS} weeks"),
        ("Format", "Hybrid (in-person + live Zoom + recorded)"),
        ("Tuition", TUITION),
        ("Prerequisites", PREREQ),
        ("Lead instructor", INSTRUCTOR),
        ("Team size", "4-6 engineers per team"),
        ("Cohort size", "Capped at 60 to preserve hands-on time"),
    ]
    add_table(doc, ["Field", "Value"], info, col_widths=[1.6, 5.4])

    add_h1(doc, "Course Description")
    add_body(doc,
        "A twelve-week intensive bootcamp on LLMs and enterprise RAG. The "
        "curriculum starts at the foundations (transformer attention, "
        "embeddings, ANN search) and moves through retrieval engineering, "
        "multimodal RAG, graph-based retrieval, production guardrails, "
        "agentic RAG, fine-tuning, and rigorous evaluation. Teams of 4-6 "
        "work on weekly projects in dedicated rooms with full multimedia "
        "setup and direct access to the Proxiant Datacenter.")

    add_h1(doc, "Weekly Cadence")
    cadence = [
        ("Saturday 11 AM - 1 PM PST", "Morning theory session"),
        ("Saturday 1 - 1:30 PM PST", "Lunch (served on-site)"),
        ("Saturday 1:30 - 4 PM PST", "Afternoon labs and exercises"),
        ("Saturday 4 - 5 PM PST", "Project presentations"),
        ("Monday 7 - 10 PM PST", "Guided lab session"),
        ("Wednesday 7 - 10 PM PST", "Guided lab session"),
        ("Tuesday 8:30 - 10 AM PST", "Summary and weekly quiz"),
    ]
    add_table(doc, ["When", "What"], cadence, col_widths=[2.4, 4.6])

    add_h1(doc, "Grading and Assessment")
    grading = [
        ("Weekly quizzes (12)", "15%", "Multiple choice + short answer, theory and labs"),
        ("Weekly team projects (12)", "30%", "Take-home, peer reviewed against a rubric"),
        ("Lab participation", "15%", "Attendance and submitted artifacts"),
        ("Paper presentations", "10%", "One per student over the cohort"),
        ("Capstone project", "25%", "Code, presentation, peer review"),
        ("Final certification exam", "5%", "PCAP-RAG written exam"),
    ]
    add_table(doc, ["Component", "Weight", "Description"], grading, col_widths=[1.8, 0.8, 4.2])
    add_body(doc, "Passing grade is 70%. PCAP-RAG certificate issued on pass.")

    add_h1(doc, "Required Tooling")
    tooling = [
        "Python 3.11+, conda or uv, PyTorch 2.4+, Hugging Face transformers and TRL",
        "sentence-transformers, FAISS, Pyserini (BM25), RAGatouille (ColBERT)",
        "DSPy, GEPA reference, TextGrad",
        "Neo4j or Memgraph, RAPTOR reference impl, GraphRAG, LightRAG",
        "open_clip, Lavis (BLIP-2)",
        "Ragas, DeBERTa-large-mnli, Presidio",
        "Google ADK, FastMCP",
        "Docker, kubectl (for the deployment labs)",
    ]
    for t in tooling:
        add_bullet(doc, t)

    add_page_break(doc)
    add_h1(doc, "Weekly Detail")

    for w in WEEKS:
        add_h2(doc, f"Week {w['num']}: {w['title']}")
        add_subtitle(doc, w["date"], size=11)
        add_body(doc, w["summary"])

        add_h3(doc, "Learning objectives")
        for o in w["objectives"]:
            add_bullet(doc, o)

        add_h3(doc, "Topics covered")
        for tname, tdesc in w["topics"]:
            add_bullet(doc, f"{tname}: {tdesc}")

        add_h3(doc, "Research papers")
        for ptitle, pnote in w["papers"]:
            add_bullet(doc, f"{ptitle}. {pnote}")

        add_h3(doc, "Labs this week")
        for lab in w["labs"]:
            add_bullet(doc, f"{lab['title']}: {lab['objective']}")

        add_divider(doc)

    out = os.path.join(ROOT, "Syllabus", "LLM_Enterprise_RAG_Syllabus.docx")
    doc.save(out)
    print(f"WROTE {out}")


# ------------------ class notes ------------------

def build_class_notes():
    for week in WEEKS:
        doc = new_doc()
        n = week["num"]
        add_title(doc, f"Week {n}: {week['title']}", size=22)
        add_subtitle(doc, week["tagline"], size=14)
        add_subtitle(doc, f"Main session date: {week['date']}", size=11)
        add_subtitle(doc, f"{COURSE_NAME} | Proxiant Academy | Instructor: {INSTRUCTOR}", size=11)
        add_divider(doc)

        add_h1(doc, "Session Overview")
        add_body(doc, DEEP[n]["intro"])

        add_h2(doc, "Learning Objectives")
        for o in week["objectives"]:
            add_bullet(doc, o)

        add_page_break(doc)
        add_h1(doc, "Detailed Notes")
        for sec in DEEP[n]["sections"]:
            add_h2(doc, sec["h"])
            add_body(doc, sec["body"])
            if "code" in sec:
                add_body(doc, "Reference snippet:")
                add_code(doc, sec["code"])

        add_page_break(doc)
        add_h1(doc, "Topics at a Glance")
        for tname, tdesc in week["topics"]:
            add_h3(doc, tname)
            add_body(doc, tdesc)

        add_h1(doc, "Research Papers")
        for ptitle, pnote in week["papers"]:
            add_h3(doc, ptitle)
            add_body(doc, pnote)

        add_h1(doc, "This Week's Labs")
        for lab in week["labs"]:
            add_h3(doc, lab["title"])
            add_body(doc, f"Objective: {lab['objective']}")
            add_body(doc, f"Deliverables: {lab['deliverables']}")

        add_h1(doc, "Review Questions")
        for q in DEEP[n]["review"]:
            add_bullet(doc, q)

        out = os.path.join(ROOT, "Class_Notes",
                           f"Week_{n:02d}_{slug(week)}_Notes.docx")
        doc.save(out)
        print(f"WROTE {out}")


# ------------------ slides ------------------

NAVY = RGBColor(0x1E, 0x27, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
ACCENT = RGBColor(0x2C, 0x5F, 0xF5)


def _band(slide, color, top, height, width=Inches(13.33)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _text(slide, text, left, top, width, height, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def _bullets(slide, items, left, top, width, height, size=15, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "  •  " + item
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _title_slide(prs, week):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _band(s, NAVY, 0, prs.slide_height)
    _text(s, "Proxiant Academy", Inches(0.6), Inches(0.4), Inches(8), Inches(0.5), size=14, color=WHITE)
    _text(s, f"Week {week['num']}", Inches(0.6), Inches(2.0), Inches(8), Inches(0.7), size=28, color=ACCENT, bold=True)
    _text(s, week["title"], Inches(0.6), Inches(2.8), Inches(12), Inches(1.4), size=44, color=WHITE, bold=True)
    _text(s, week["tagline"], Inches(0.6), Inches(4.6), Inches(12), Inches(1.0), size=20, color=WHITE)
    _text(s, week["date"], Inches(0.6), Inches(6.6), Inches(8), Inches(0.4), size=12, color=WHITE)
    _text(s, f"{COURSE_NAME} | Instructor: {INSTRUCTOR}", Inches(0.6), Inches(7.0), Inches(10), Inches(0.4), size=12, color=WHITE)


def _section_slide(prs, label, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _band(s, LIGHT, 0, prs.slide_height)
    _band(s, NAVY, Inches(3.0), Inches(0.05))
    _text(s, label, Inches(0.7), Inches(2.4), Inches(8), Inches(0.5), size=16, color=ACCENT, bold=True)
    _text(s, title, Inches(0.7), Inches(2.9), Inches(12), Inches(1.4), size=36, color=NAVY, bold=True)


def _content_slide(prs, header, title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _band(s, NAVY, 0, Inches(0.45))
    _text(s, header, Inches(0.5), Inches(0.08), Inches(12), Inches(0.3), size=11, color=WHITE)
    _text(s, title, Inches(0.5), Inches(0.7), Inches(12.5), Inches(0.8), size=26, color=NAVY, bold=True)
    _bullets(s, bullets, Inches(0.7), Inches(1.7), Inches(12), Inches(5.6))
    _band(s, NAVY, prs.slide_height - Inches(0.2), Inches(0.2))


def _code_slide(prs, header, title, code):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _band(s, NAVY, 0, Inches(0.45))
    _text(s, header, Inches(0.5), Inches(0.08), Inches(12), Inches(0.3), size=11, color=WHITE)
    _text(s, title, Inches(0.5), Inches(0.7), Inches(12.5), Inches(0.8), size=26, color=NAVY, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x2C)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    for i, line in enumerate(code.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(13)
        r.font.color.rgb = RGBColor(0xE6, 0xEC, 0xF5)
    _band(s, NAVY, prs.slide_height - Inches(0.2), Inches(0.2))


def _closing_slide(prs, week):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _band(s, NAVY, 0, prs.slide_height)
    _text(s, "This Week's Labs", Inches(0.7), Inches(0.9), Inches(10), Inches(0.6), size=18, color=ACCENT, bold=True)
    for i, lab in enumerate(week["labs"]):
        _text(s, f"Lab {chr(65+i)}: {lab['title']}",
              Inches(0.7), Inches(1.6 + i*0.9), Inches(12), Inches(0.5),
              size=18, color=WHITE, bold=True)
        _text(s, lab["objective"],
              Inches(0.9), Inches(2.0 + i*0.9), Inches(12), Inches(0.7),
              size=12, color=WHITE)
    _text(s, "Quiz Tuesday | Project due next Saturday | Instructor: " + INSTRUCTOR,
          Inches(0.7), Inches(7.0), Inches(12), Inches(0.4), size=12, color=WHITE)


def build_slides():
    for week in WEEKS:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        _title_slide(prs, week)

        _section_slide(prs, "Section 1", "Today's Map")
        _content_slide(prs, f"Week {week['num']} | {week['title']}",
                       "Where we are in the bootcamp", [week["summary"][:380]])
        _content_slide(prs, f"Week {week['num']} | {week['title']}",
                       "Learning Objectives", week["objectives"])

        _section_slide(prs, "Section 2", "Concepts")
        for tname, tdesc in week["topics"]:
            _content_slide(prs, f"Week {week['num']} | {week['title']}", tname, [tdesc])

        deep = DEEP[week["num"]]
        _section_slide(prs, "Section 3", "Deeper Look")
        for sec in deep["sections"]:
            if "code" in sec:
                _code_slide(prs, f"Week {week['num']} | {week['title']}", sec["h"], sec["code"])
            else:
                _content_slide(prs, f"Week {week['num']} | {week['title']}", sec["h"], [sec["body"]])

        _section_slide(prs, "Section 4", "Papers")
        for ptitle, pnote in week["papers"]:
            _content_slide(prs, f"Week {week['num']} | {week['title']}", ptitle, [pnote])

        _section_slide(prs, "Section 5", "Review")
        _content_slide(prs, f"Week {week['num']} | {week['title']}",
                       "Review Questions", deep["review"])

        _closing_slide(prs, week)

        out = os.path.join(ROOT, "Slides",
                           f"Week_{week['num']:02d}_{slug(week)}.pptx")
        prs.save(out)
        print(f"WROTE {out}")


# ------------------ lab walkthroughs ------------------

def _safe(t):
    out = []
    for c in t.lower():
        if c.isalnum():
            out.append(c)
        elif c in " -_":
            out.append("_")
    return "".join(out)[:60].strip("_")


def build_labs():
    for week in WEEKS:
        n = week["num"]
        for i, lab in enumerate(week["labs"]):
            doc = new_doc()
            add_title(doc, f"Week {n} | Lab {chr(65+i)}", size=22)
            add_subtitle(doc, lab["title"], size=14)
            sess = "Monday" if i == 0 else "Wednesday" if i == 1 else "Friday"
            add_subtitle(doc, f"Bootcamp week: {week['title']} | Session: {sess} | Instructor: {INSTRUCTOR}", size=11)
            add_divider(doc)

            add_h1(doc, "Lab Objective")
            add_body(doc, lab["objective"])

            add_h1(doc, "Prerequisites")
            add_body(doc, lab["prereqs"])

            add_h1(doc, "Environment Setup")
            add_body(doc, "Cluster workspace from week 1 plus week-specific dependencies in requirements.txt that ships with the lab repo.")

            add_h1(doc, "Walkthrough")
            for step in lab["steps"]:
                add_numbered(doc, step)

            add_h1(doc, "Deliverables")
            add_body(doc, lab["deliverables"])

            add_h1(doc, "Grading Rubric")
            rubric = [
                ("Code correctness", "30", "Code runs as specified; results in expected ranges"),
                ("Engineering quality", "20", "Readable, no dead code, sensible names"),
                ("Observability", "15", "Logs, traces, or charts as required"),
                ("Analysis", "20", "Written reflection identifies real tradeoffs"),
                ("Reproducibility", "15", "Peer can re-run by following the README"),
            ]
            add_table(doc, ["Criterion", "Points", "Description"], rubric, col_widths=[2.0, 0.8, 4.2])

            add_h1(doc, "Common Pitfalls")
            add_bullet(doc, "Skipping the verification checks. Each one catches a frequent error.")
            add_bullet(doc, "Hard-coding paths. Use the lab repo's config.")
            add_bullet(doc, "Skipping the reflection. The TA grading uses it heavily.")

            fname = f"Week_{n:02d}_Lab_{chr(65+i)}_{_safe(lab['title'])}.docx"
            out = os.path.join(ROOT, "Lab_Walkthroughs", fname)
            doc.save(out)
            print(f"WROTE {out}")


# ------------------ quizzes ------------------

def build_quizzes():
    for week in WEEKS:
        for with_answers in (False, True):
            n = week["num"]
            doc = new_doc()
            qz = QUIZZES[n]
            suffix = " — Answer Key" if with_answers else ""
            add_title(doc, f"Week {n} Quiz{suffix}", size=22)
            add_subtitle(doc, week["title"], size=14)
            add_subtitle(doc, f"Tuesday 8:30 to 10:00 AM PST | 30 minutes | 60 points | Instructor: {INSTRUCTOR}", size=11)
            add_divider(doc)

            add_h1(doc, "Instructions")
            add_body(doc, "12 multiple-choice questions (3 points each) and 3 short-answer "
                          "questions (8 points each).")

            add_h1(doc, "Part A: Multiple Choice")
            for i, (q, options, ans, expl) in enumerate(qz["mcq"], 1):
                add_h2(doc, f"Q{i}. {q}")
                for j, opt in enumerate(options):
                    text = f"{LETTERS[j]}. {opt}"
                    if with_answers and j == ans:
                        text += "   [correct]"
                    add_bullet(doc, text)
                if with_answers:
                    add_body(doc, f"Explanation: {expl}")

            add_h1(doc, "Part B: Short Answer")
            for i, sq in enumerate(qz["short"], 1):
                add_h2(doc, f"S{i}. {sq}")
                if not with_answers:
                    add_body(doc, "(Write 3 to 6 sentences in the space below.)")
                else:
                    add_body(doc, "Model answers reviewed live in the Tuesday quiz session.")

            fname = f"Week_{n:02d}_Quiz{'_Solution' if with_answers else ''}.docx"
            out = os.path.join(ROOT, "Weekly_Quizzes", fname)
            doc.save(out)
            print(f"WROTE {out}")


# ------------------ projects ------------------

def build_projects():
    for week in WEEKS:
        for solution in (False, True):
            n = week["num"]
            proj = PROJECTS[n]
            doc = new_doc()
            suffix = " — Reference Solution Sketch" if solution else ""
            add_title(doc, f"Week {n} Project Brief{suffix}", size=22)
            add_subtitle(doc, proj["title"], size=14)
            add_subtitle(doc, f"Bootcamp theme: {week['title']} | Due next Saturday | Instructor: {INSTRUCTOR}", size=11)
            add_divider(doc)

            add_h1(doc, "Project Summary")
            add_body(doc, proj["summary"])

            add_h1(doc, "Required Tools and Libraries")
            for t in proj["tools"]:
                add_bullet(doc, t)

            add_h1(doc, "Deliverables")
            for d in proj["deliverables"]:
                add_bullet(doc, d)

            add_h1(doc, "Grading Rubric")
            rubric = [
                ("Correctness", "30", "Code runs as specified; results meet the bar"),
                ("Engineering quality", "20", "Clean, readable, tested where it matters"),
                ("Evaluation rigor", "20", "Eval design is honest; numbers are reproducible"),
                ("Observability and reproducibility", "15", "README plus traces let a peer rerun in under 30 minutes"),
                ("Written reflection", "15", "Tradeoffs explicit, failure cases owned"),
            ]
            add_table(doc, ["Criterion", "Points", "Description"], rubric, col_widths=[2.4, 0.8, 4.0])
            add_body(doc, proj["rubric_extra"])

            add_h1(doc, "Submission")
            add_bullet(doc, "Push to a private GitHub repo and add the TAs as collaborators.")
            add_bullet(doc, "Tag the submission commit 'week-XX-submit'.")
            add_bullet(doc, "Post the repo link in the bootcamp Discord by Saturday 10 AM PST.")

            if solution:
                add_h1(doc, "Reference Solution Sketch")
                add_body(doc, "A clean reference solution emphasizes honest evaluation over polished UI. Skeletons for each deliverable are in the lab repo branch 'reference-week-XX'.")

                add_h2(doc, "Architecture")
                add_body(doc, f"Built around the patterns introduced in week {n}'s class notes. Follow the recommended defaults; deviate only with measured justification.")

                add_h2(doc, "Common Failure Modes")
                add_bullet(doc, "Skipping evaluation rigor in favor of better-looking implementation.")
                add_bullet(doc, "Hyperparameter changes without ablation.")
                add_bullet(doc, "Hidden state across runs (no seed, no version pin).")

            fname = f"Week_{n:02d}_Project{'_Solution' if solution else ''}.docx"
            out = os.path.join(ROOT, "Weekly_Projects", fname)
            doc.save(out)
            print(f"WROTE {out}")


# ------------------ research paper guides ------------------

def build_papers():
    for week in WEEKS:
        n = week["num"]
        doc = new_doc()
        add_title(doc, f"Week {n} Research Paper Reading Guide", size=22)
        add_subtitle(doc, week["title"], size=14)
        add_subtitle(doc, f"Read before Saturday's main session | Instructor: {INSTRUCTOR}", size=11)
        add_divider(doc)

        add_h1(doc, "Framing for This Week")
        add_body(doc,
            "Read in the order presented. The first paper establishes the "
            "foundational concept; subsequent papers add depth or contrast.")

        add_h1(doc, "Papers")
        for ptitle, pnote in week["papers"]:
            add_h2(doc, ptitle)
            add_body(doc, pnote)
            add_h3(doc, "What to look for")
            add_bullet(doc, "Main claim in one sentence.")
            add_bullet(doc, "The experimental setup that supports the claim.")
            add_bullet(doc, "The one assumption you would challenge.")
            add_bullet(doc, "How the paper relates to this week's lab work.")

        add_h1(doc, "Discussion Questions")
        add_bullet(doc, "Which of these papers will still be relevant in 5 years? Which will not?")
        add_bullet(doc, "Pick one experimental setup choice the authors made. Suggest a better alternative.")
        add_bullet(doc, "What is the most expensive part of reproducing these results?")
        add_bullet(doc, "Connect the methods to a problem your team is solving in this week's project.")

        add_h1(doc, "Presenter Notes")
        add_body(doc,
            "The assigned presenter leads a 15-minute walkthrough of one "
            "paper, followed by 20 minutes of group discussion. Strong "
            "opinions welcome, evidence required.")

        add_h1(doc, "Submission")
        add_bullet(doc, "Each student submits a 200-word reflection on Discord by Friday EOD.")
        add_bullet(doc, "Reflections feed into the Saturday discussion and the Tuesday quiz.")

        fname = f"Week_{n:02d}_Paper_Reading_{slug(week)}.docx"
        out = os.path.join(ROOT, "Research_Papers", fname)
        doc.save(out)
        print(f"WROTE {out}")


# ------------------ certification exam ------------------

SCENARIOS = [
    {
        "title": "Scenario A: Chunking strategy selection",
        "prompt": (
            "Your team is building a RAG over a corpus of 50K policy "
            "documents, each 5 to 30 pages. Queries are mostly multi-paragraph "
            "answers requiring context from multiple sections. P95 latency "
            "budget: 3 seconds. Pick a chunking strategy and defend the "
            "choice. Include cost estimates."
        ),
        "rubric": [
            "Recognizes long-document, multi-section queries need hierarchical or contextual chunking.",
            "Estimates ingestion cost of contextual chunking honestly.",
            "Considers latency budget against chunking choice.",
            "Names a specific embedder appropriate to the choice.",
            "Plans an eval to validate the choice rather than asserting it.",
        ],
    },
    {
        "title": "Scenario B: Retrieval funnel design",
        "prompt": (
            "Design a four-stage retrieval funnel for a 10M-document corpus "
            "with a 200ms retrieval budget. Specify each stage, candidate "
            "counts, and which components you would skip if forced to halve "
            "the latency."
        ),
        "rubric": [
            "Sparse + dense at stage 1 with realistic candidate counts.",
            "Matryoshka prune with concrete dimensions (e.g., 64 for prune).",
            "ColBERT and cross-encoder placement justified by cost.",
            "Sensible cut-down for 100ms budget (typically drop ColBERT).",
            "Mentions latency budgeting per stage.",
        ],
    },
    {
        "title": "Scenario C: Production guardrails for a healthcare RAG",
        "prompt": (
            "Build the guardrail and grounding pipeline for a healthcare RAG "
            "that answers questions about clinical guidelines. Must comply "
            "with HIPAA. Must hit a 4-second P95. Document the threat model "
            "and the pipeline."
        ),
        "rubric": [
            "PII detection and redaction with explicit policy.",
            "Prompt injection defense for untrusted retrieved content.",
            "NLI-based grounding verification.",
            "Citation requirement for clinical claims.",
            "Latency budget allocation across stages.",
        ],
    },
    {
        "title": "Scenario D: Evaluation suite design",
        "prompt": (
            "Design an evaluation suite for a multi-hop legal RAG. Include "
            "retrieval metrics, generation metrics, robustness tests, and "
            "an LLM-as-judge component. Address each known bias and explain "
            "your mitigation."
        ),
        "rubric": [
            "Right retrieval metric (NDCG with graded relevance for legal).",
            "Ragas faithfulness and context coverage.",
            "RGB-style robustness tests, named explicitly.",
            "LLM judge with position randomization, length normalization, different judge model.",
            "Reproducibility plan with seed and version pinning.",
        ],
    },
]


def build_cert_test():
    for with_answers in (False, True):
        doc = new_doc()
        suffix = " — Answer Key" if with_answers else ""
        add_title(doc, f"PCAP-RAG Final Certification Exam{suffix}", size=22)
        add_subtitle(doc, COURSE_NAME, size=14)
        add_subtitle(doc, f"Duration: 3 hours | Total: 240 points | Pass: 168 (70%) | Instructor: {INSTRUCTOR}", size=11)
        add_divider(doc)

        add_h1(doc, "Exam Structure")
        structure = [
            ("Part A", "60 multiple-choice questions", "3 points each = 180 points"),
            ("Part B", "4 scenario-based long-form responses", "15 points each = 60 points"),
            ("Total", "", "240 points"),
        ]
        add_table(doc, ["Part", "Questions", "Scoring"], structure, col_widths=[1.0, 3.8, 2.0])

        add_h1(doc, "Instructions")
        add_bullet(doc, "Closed-notes for Part A. Open-notes for Part B (course materials only).")
        add_bullet(doc, "Allocate roughly 90 minutes for Part A, 90 minutes for Part B.")
        add_bullet(doc, "Partial credit awarded on Part B based on the rubric for each scenario.")

        add_h1(doc, "Part A: Multiple Choice (60 questions)")
        idx = 1
        for week_num in range(1, 13):
            mcqs = QUIZZES[week_num]["mcq"][:5]
            for (q, options, ans, expl) in mcqs:
                add_h2(doc, f"Q{idx} (Week {week_num}). {q}")
                for j, opt in enumerate(options):
                    text = f"{LETTERS[j]}. {opt}"
                    if with_answers and j == ans:
                        text += "   [correct]"
                    add_bullet(doc, text)
                if with_answers:
                    add_body(doc, f"Explanation: {expl}")
                idx += 1

        add_h1(doc, "Part B: Scenarios (4 questions)")
        for i, s in enumerate(SCENARIOS, 1):
            add_h2(doc, f"Scenario {i}: {s['title']}")
            add_body(doc, s["prompt"])
            if with_answers:
                add_body(doc, "Grading rubric (3 points per item):")
                for item in s["rubric"]:
                    add_bullet(doc, item)

        add_h1(doc, "Certification")
        add_body(doc,
            "A passing grade earns the PCAP-RAG certificate, recognized "
            "across the Proxiant alumni network and partner employers. The "
            "certificate is valid for 3 years; renewal requires either a "
            "continuing-education credit or re-examination.")

        fname = f"PCAP_RAG_Final_Exam{'_Solution' if with_answers else ''}.docx"
        out = os.path.join(ROOT, "Certification_Test", fname)
        doc.save(out)
        print(f"WROTE {out}")


if __name__ == "__main__":
    build_catalogue()
    build_syllabus()
    build_class_notes()
    build_labs()
    build_quizzes()
    build_projects()
    build_papers()
    build_cert_test()
    build_slides()
