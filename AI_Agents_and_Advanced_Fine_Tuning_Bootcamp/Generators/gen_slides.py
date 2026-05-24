"""Generate PPTX slide decks: one per week, ~16-20 slides each."""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from bootcamp_content import WEEKS, slug
from class_notes_deep import DEEP

ROOT = "/Users/pkr465/work/AI-Course-Material/AI_Agents_and_Advanced_Fine_Tuning_Bootcamp"

NAVY = RGBColor(0x1E, 0x27, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
ACCENT = RGBColor(0x2C, 0x5F, 0xF5)


def add_band(slide, color, top, height, width=Inches(13.33)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_bullets(slide, items, left, top, width, height, size=16, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "  •  " + item
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color


def title_slide(prs, week):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_band(s, NAVY, 0, prs.slide_height)
    add_text(s, "Proxiant Academy", Inches(0.6), Inches(0.4),
             Inches(8), Inches(0.5), size=14, color=WHITE)
    add_text(s, f"Week {week['num']}", Inches(0.6), Inches(2.0),
             Inches(8), Inches(0.7), size=28, color=ACCENT, bold=True)
    add_text(s, week["title"], Inches(0.6), Inches(2.8),
             Inches(12), Inches(1.4), size=44, color=WHITE, bold=True)
    add_text(s, week["tagline"], Inches(0.6), Inches(4.6),
             Inches(12), Inches(1.0), size=20, color=WHITE)
    add_text(s, week["date"], Inches(0.6), Inches(6.6),
             Inches(8), Inches(0.4), size=12, color=WHITE)
    add_text(s, "AI Agents and Advanced Fine-Tuning Bootcamp",
             Inches(0.6), Inches(7.0), Inches(10), Inches(0.4),
             size=12, color=WHITE)


def section_slide(prs, label, title):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_band(s, LIGHT, 0, prs.slide_height)
    add_band(s, NAVY, Inches(3.0), Inches(0.05))
    add_text(s, label, Inches(0.7), Inches(2.4), Inches(8), Inches(0.5),
             size=16, color=ACCENT, bold=True)
    add_text(s, title, Inches(0.7), Inches(2.9), Inches(12), Inches(1.4),
             size=36, color=NAVY, bold=True)


def content_slide(prs, header, title, bullets, body=None):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_band(s, NAVY, 0, Inches(0.45))
    add_text(s, header, Inches(0.5), Inches(0.08), Inches(12), Inches(0.3),
             size=11, color=WHITE)
    add_text(s, title, Inches(0.5), Inches(0.7), Inches(12.5), Inches(0.8),
             size=26, color=NAVY, bold=True)
    if body:
        add_text(s, body, Inches(0.5), Inches(1.6), Inches(12.5), Inches(1.2),
                 size=14, color=DARK)
        top = Inches(2.9)
        h = Inches(4.5)
    else:
        top = Inches(1.7)
        h = Inches(5.6)
    add_bullets(s, bullets, Inches(0.7), top, Inches(12), h, size=15)
    add_band(s, NAVY, prs.slide_height - Inches(0.2), Inches(0.2))


def code_slide(prs, header, title, code):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_band(s, NAVY, 0, Inches(0.45))
    add_text(s, header, Inches(0.5), Inches(0.08), Inches(12), Inches(0.3),
             size=11, color=WHITE)
    add_text(s, title, Inches(0.5), Inches(0.7), Inches(12.5), Inches(0.8),
             size=26, color=NAVY, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.5), Inches(1.7),
                             Inches(12.3), Inches(5.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x2C)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    for i, line in enumerate(code.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(13)
        r.font.color.rgb = RGBColor(0xE6, 0xEC, 0xF5)
    add_band(s, NAVY, prs.slide_height - Inches(0.2), Inches(0.2))


def closing_slide(prs, week):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_band(s, NAVY, 0, prs.slide_height)
    add_text(s, "This Week's Labs", Inches(0.7), Inches(0.9),
             Inches(10), Inches(0.6), size=18, color=ACCENT, bold=True)
    for i, lab in enumerate(week["labs"]):
        add_text(s, f"Lab {chr(65+i)}: {lab['title']}",
                 Inches(0.7), Inches(1.6 + i*0.9), Inches(12), Inches(0.5),
                 size=18, color=WHITE, bold=True)
        add_text(s, lab["objective"],
                 Inches(0.9), Inches(2.0 + i*0.9), Inches(12), Inches(0.7),
                 size=12, color=WHITE)
    add_text(s, "Quiz on Wednesday | Project due next Sunday",
             Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
             size=12, color=WHITE)


def build_one(week):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_slide(prs, week)

    section_slide(prs, "Section 1", "Today's Map")
    content_slide(prs,
                  f"Week {week['num']} | {week['title']}",
                  "Where we are in the bootcamp",
                  [week["summary"][:380]],
                  body=None)
    content_slide(prs,
                  f"Week {week['num']} | {week['title']}",
                  "Learning Objectives",
                  week["objectives"])

    section_slide(prs, "Section 2", "Concepts")
    for tname, tdesc in week["topics"]:
        content_slide(prs,
                      f"Week {week['num']} | {week['title']}",
                      tname,
                      [tdesc])

    deep = DEEP[week["num"]]
    section_slide(prs, "Section 3", "Deeper Look")
    for sec in deep["sections"]:
        if "code" in sec:
            code_slide(prs,
                       f"Week {week['num']} | {week['title']}",
                       sec["h"],
                       sec["code"])
        else:
            content_slide(prs,
                          f"Week {week['num']} | {week['title']}",
                          sec["h"],
                          [sec["body"]])

    section_slide(prs, "Section 4", "Papers")
    for ptitle, pnote in week["papers"]:
        content_slide(prs,
                      f"Week {week['num']} | {week['title']}",
                      ptitle,
                      [pnote])

    section_slide(prs, "Section 5", "Review")
    content_slide(prs,
                  f"Week {week['num']} | {week['title']}",
                  "Review Questions",
                  deep["review"])

    closing_slide(prs, week)

    out = os.path.join(ROOT, "Slides",
                       f"Week_{week['num']:02d}_{slug(week)}.pptx")
    prs.save(out)
    print(f"WROTE {out}")


if __name__ == "__main__":
    for w in WEEKS:
        build_one(w)
